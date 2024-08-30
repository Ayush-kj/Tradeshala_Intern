from pptx import Presentation
from pptx.util import Inches

# Create a presentation object
prs = Presentation()

# Title Slide
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Market Information for Asset Classes"
subtitle.text = "Major Currencies, Cross Currencies, Exotic Currencies, and Commodities"

# Slide 1 - Major Currencies
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Major Currencies"
content.text = ("1. EUR/USD\n2. USD/JPY\n3. GBP/USD\n4. USD/CHF\n5. AUD/USD")

# Slide 2 - Cross Currencies
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Cross Currencies"
content.text = ("1. EUR/GBP\n2. EUR/JPY\n3. GBP/JPY\n4. AUD/NZD\n5. CHF/JPY")

# Slide 3 - Exotic Currencies
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Exotic Currencies"
content.text = ("1. USD/INR\n2. USD/TRY\n3. USD/ZAR")

# Slide 4 - Commodities
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Commodities"
content.text = ("1. Gold (XAU/USD)\n2. Crude Oil (WTI/USD)")

# Slide 5 - Major Currencies Market Information
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Major Currencies Market Information"

# Adding a table
rows = 6
cols = 6
table_data = [
    ["Pair", "1 Base = X Quote", "Bid Price", "Ask Price", "Spread (Mini Lot)", "Spread (Standard Lot)"],
    ["EUR/USD", "1.08436", "1.08430", "1.08436", "$0.60", "$6.00"],
    ["USD/JPY", "149.78", "149.77", "149.78", "$0.10", "$1.00"],
    ["GBP/USD", "1.24567", "1.24562", "1.24567", "$0.50", "$5.00"],
    ["USD/CHF", "0.90090", "0.90085", "0.90090", "$0.50", "$5.00"],
    ["AUD/USD", "0.66543", "0.66538", "0.66543", "$0.50", "$5.00"]
]

table = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(1.5), Inches(9), Inches(2.5)).table

# Set column widths
for i in range(cols):
    table.columns[i].width = Inches(1.5)

# Add data to table
for i in range(rows):
    for j in range(cols):
        table.cell(i, j).text = table_data[i][j]

# Slide 6 - Interpretation of EUR/USD
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Interpretation of EUR/USD Two-Way Market"
content.text = ("Base Currency: EUR\nQuote Currency: USD\n"
                "If you want to buy EUR/USD, you pay 1.08436 USD for 1 EUR (ask price).\n"
                "If you want to sell, you get 1.08430 USD for 1 EUR (bid price).\n"
                "The spread is 0.6 pips.")

# Slide 7 - Interpretation of USD/JPY
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Interpretation of USD/JPY Two-Way Market"
content.text = ("Base Currency: USD\nQuote Currency: JPY\n"
                "Buying 1 USD costs 149.78 JPY (ask price).\n"
                "Selling 1 USD gets you 149.77 JPY (bid price).\n"
                "The spread is 1 pip.")

# Slide 8 - Interpretation of XAU/USD
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Interpretation of XAU/USD Two-Way Market"
content.text = ("Base Currency: XAU (Gold)\nQuote Currency: USD\n"
                "Buying 1 ounce of gold costs 2362.04 USD (ask price).\n"
                "Selling 1 ounce gets you 2361.04 USD (bid price).\n"
                "The spread is 1 pip.")

# Slide 9 - Profit and Loss Calculations
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Profit and Loss Calculations"
content.text = ("1. EUR/USD (1 PIP movement): 1.08426 to 1.08436\n"
                "Units: 100,000 (Standard Lot)\n"
                "Profit: $10\n\n"
                "2. USD/JPY (1 PIP movement): 149.77 to 149.78\n"
                "Units: 10,000 (Mini Lot)\n"
                "Profit: ¥100 = $0.67\n\n"
                "3. XAU/USD (1 Tick movement): 2361.04 to 2362.04\n"
                "Units: 10,000 (Mini Lot)\n"
                "Profit: $10,000\n\n"
                "4. BTC/USD (1 Pip movement): 67,544 to 67,545\n"
                "Units: 10,000 (Mini Lot)\n"
                "Profit: $10,000")

# Slide 10 - Conclusion
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Conclusion"
content.text = "Summary of findings and interpretations."

# Save the presentation
prs.save('Market_Information_Presentation.pptx')
